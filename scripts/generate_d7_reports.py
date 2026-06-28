"""
generate_d7_reports.py
======================
Bluestock Mutual Fund Capstone – Deliverable D7: Final Documentation

Generates:
  reports/Final_Report.docx
  reports/Final_Report.pdf
  reports/Presentation.pptx

Run with:
    python scripts/generate_d7_reports.py
"""

from __future__ import annotations

import os
import sys
import io
from pathlib import Path
from datetime import date

# ── Path setup ────────────────────────────────────────────────────────────────
PROJECT_ROOT = Path(__file__).resolve().parent.parent
REPORTS_DIR  = PROJECT_ROOT / "reports"
DOCS_DIR     = PROJECT_ROOT / "docs"
DATA_PROC    = PROJECT_ROOT / "data" / "processed"

REPORTS_DIR.mkdir(exist_ok=True)

# ── Image helpers ─────────────────────────────────────────────────────────────
SCREENSHOT_DIR = DOCS_DIR / "dashboard_screenshots"
CHART_IMAGES   = {
    "nav_trends"           : DOCS_DIR / "01_nav_trends.png",
    "sip_inflows"          : DOCS_DIR / "02_sip_inflows.png",
    "aum_growth"           : DOCS_DIR / "03_aum_growth.png",
    "risk_return"          : DOCS_DIR / "04_risk_return.png",
    "sector_alloc"         : DOCS_DIR / "05_sector_allocations.png",
    "txn_correlation"      : DOCS_DIR / "06_transaction_correlation.png",
    "category_inflows"     : DOCS_DIR / "07_category_comparison.png",
    "fund_vs_bench"        : DOCS_DIR / "fund_vs_benchmark.png",
    "dash_overview"        : SCREENSHOT_DIR / "1_executive_overview.png",
    "dash_perf"            : SCREENSHOT_DIR / "2_performance_analytics.png",
    "dash_portfolio"       : SCREENSHOT_DIR / "3_portfolio_allocation.png",
    "dash_risk"            : SCREENSHOT_DIR / "4_risk_analytics.png",
    "var_comparison"       : DATA_PROC / "chart_01_var_comparison.png",
    "return_dist"          : DATA_PROC / "chart_02_return_distribution.png",
    "category_cohort"      : DATA_PROC / "chart_03_category_cohort.png",
    "amc_cohort"           : DATA_PROC / "chart_04_amc_cohort.png",
    "correlation"          : DATA_PROC / "chart_06_correlation_heatmap.png",
    "rec_scores"           : DATA_PROC / "chart_08_recommender_scores.png",
    "risk_return_adv"      : DATA_PROC / "chart_09_risk_return.png",
    "sharpe_violin"        : DATA_PROC / "chart_11_sharpe_violin.png",
    "max_drawdown"         : DATA_PROC / "chart_12_max_drawdown.png",
}

TODAY = date.today().strftime("%d %B %Y")

# ─────────────────────────────────────────────────────────────────────────────
# SECTION 1: DOCX REPORT
# ─────────────────────────────────────────────────────────────────────────────

def build_docx():
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = Document()

    # ── Page margins ───────────────────────────────────────────────
    for section in doc.sections:
        section.top_margin    = Cm(2.0)
        section.bottom_margin = Cm(2.0)
        section.left_margin   = Cm(2.5)
        section.right_margin  = Cm(2.5)

    # ── Colour palette ─────────────────────────────────────────────
    NAVY    = RGBColor(0x0D, 0x1B, 0x2A)
    BLUE    = RGBColor(0x1B, 0x6C, 0xA8)
    ACCENT  = RGBColor(0x2E, 0xCC, 0x71)
    GREY    = RGBColor(0x55, 0x60, 0x6B)
    WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

    def set_run_color(run, rgb):
        run.font.color.rgb = rgb

    def add_heading(text, level=1, color=NAVY):
        h = doc.add_heading(text, level=level)
        for run in h.runs:
            run.font.color.rgb = color
        return h

    def add_body(text, bold=False, color=None, size=11):
        p = doc.add_paragraph()
        run = p.add_run(text)
        run.bold = bold
        run.font.size = Pt(size)
        if color:
            run.font.color.rgb = color
        return p

    def add_bullet(text, level=0):
        p = doc.add_paragraph(style="List Bullet")
        p.paragraph_format.left_indent = Cm(level * 0.5 + 0.5)
        run = p.add_run(text)
        run.font.size = Pt(10.5)
        return p

    def add_img(path, width=Inches(5.5), caption=None):
        p = str(path)
        if not Path(p).exists():
            doc.add_paragraph(f"[Chart not found: {Path(p).name}]")
            return
        doc.add_picture(p, width=width)
        if caption:
            cp = doc.add_paragraph(caption)
            cp.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for run in cp.runs:
                run.font.size = Pt(9)
                run.font.color.rgb = GREY
                run.italic = True

    def add_table_row(table, cells, bold=False, bg=None):
        row = table.add_row()
        for i, txt in enumerate(cells):
            cell = row.cells[i]
            cell.text = str(txt)
            for run in cell.paragraphs[0].runs:
                run.bold = bold
                run.font.size = Pt(9)

    # ──────────────────────────────────────────────────────────────
    # TITLE PAGE
    # ──────────────────────────────────────────────────────────────
    doc.add_paragraph()
    t = doc.add_heading("Bluestock Mutual Fund Analytics Platform", 0)
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in t.runs:
        run.font.color.rgb = NAVY
        run.font.size = Pt(24)

    sub = doc.add_paragraph("Final Project Report – Capstone Submission")
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in sub.runs:
        run.font.color.rgb = BLUE
        run.font.size = Pt(14)

    meta_lines = [
        f"Date: {TODAY}",
        "Author: Bluestock Internship Team",
        "Institution: Bluestock Fintech Private Limited",
        "Project: Mutual Fund Analytics Capstone (D1 – D7)",
    ]
    for line in meta_lines:
        p = doc.add_paragraph(line)
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in p.runs:
            run.font.size = Pt(11)
            run.font.color.rgb = GREY

    doc.add_page_break()

    # ──────────────────────────────────────────────────────────────
    # 1. EXECUTIVE SUMMARY
    # ──────────────────────────────────────────────────────────────
    add_heading("1. Executive Summary", 1, NAVY)
    add_body(
        "The Bluestock Mutual Fund Analytics Platform is an enterprise-grade financial "
        "data engineering and analytics solution. Built as a seven-deliverable capstone "
        "project, the platform ingests 10 raw CSV datasets containing over 46,000 NAV "
        "records, 32,778 investor transactions, and active portfolio holdings. It "
        "orchestrates a full ETL pipeline, populates a relational SQLite star-schema "
        "database, and surfaces advanced financial analytics through an interactive "
        "Streamlit dashboard and a rule-based recommendation engine."
    )
    doc.add_paragraph()
    add_body("Key Platform Metrics:", bold=True)
    metrics = [
        "Total Platform AUM: INR 62.74 Lakh Crore across 10 fund houses",
        "Monthly Industry SIP Inflows: INR 31,002 Crore (Dec 2025)",
        "Active SIP Accounts: 9.35 Crore",
        "Unique Investor Transactions: 32,778",
        "NAV Records: 46,000+ across 40 schemes (2022–2025)",
        "Database Tables: 11 (2 dimension + 9 fact/helper)",
        "Analytical Charts Generated: 19+ (static + interactive)",
        "Deliverables Completed: D1 ETL, D2 Database, D3 EDA, D4 Performance, D5 Dashboard, D6 Advanced Analytics, D7 Documentation",
    ]
    for m in metrics:
        add_bullet(m)

    # ──────────────────────────────────────────────────────────────
    # 2. PROBLEM STATEMENT
    # ──────────────────────────────────────────────────────────────
    doc.add_page_break()
    add_heading("2. Problem Statement", 1, NAVY)
    add_body(
        "India's mutual fund industry manages over INR 50 Lakh Crore in assets across "
        "thousands of schemes, yet retail investors often lack access to professional-grade "
        "risk analytics, portfolio diversification tools, and personalised investment guidance. "
        "Existing platforms focus on raw return figures without adequately communicating "
        "risk-adjusted performance, drawdown exposure, or correlation-driven diversification."
    )
    doc.add_paragraph()
    add_body("Core Problems Addressed:", bold=True)
    problems = [
        "Fragmented raw data across multiple CSVs with missing values and inconsistent formats",
        "No centralised relational database for multi-dimensional querying",
        "Retail investors unable to compare funds on Sharpe Ratio, VaR, or drawdown metrics",
        "Absence of a personalised, rule-based fund recommendation engine",
        "No interactive dashboard for real-time exploration of fund performance",
    ]
    for p in problems:
        add_bullet(p)

    # ──────────────────────────────────────────────────────────────
    # 3. OBJECTIVES
    # ──────────────────────────────────────────────────────────────
    doc.add_page_break()
    add_heading("3. Objectives", 1, NAVY)
    objectives = [
        ("D1 – ETL Pipeline", "Ingest, clean, and standardise 10 raw CSV datasets into a validated processed layer."),
        ("D2 – SQLite Database", "Design and populate a star-schema relational database with zero foreign-key violations."),
        ("D3 – Exploratory Data Analysis", "Generate 7+ premium static charts revealing key trends across NAV, AUM, SIP, and transactions."),
        ("D4 – Performance Analytics", "Compute CAGR, Sharpe, Sortino, Beta, Alpha, Max Drawdown, and Rolling Returns for all 40 schemes."),
        ("D5 – Streamlit Dashboard", "Build a multi-page interactive dashboard with live filtering, Bollinger Bands, wealth simulators, and PDF export."),
        ("D6 – Advanced Analytics", "Implement VaR analysis, cohort analysis, correlation heatmaps, and a rule-based recommendation engine."),
        ("D7 – Final Documentation", "Produce a professional report (DOCX + PDF) and presentation (PPTX) summarising all deliverables."),
    ]
    for title, desc in objectives:
        add_body(f"{title}: ", bold=True, size=11)
        p = doc.paragraphs[-1]
        run2 = p.add_run(desc)
        run2.font.size = Pt(11)

    # ──────────────────────────────────────────────────────────────
    # 4. DATASET DESCRIPTION
    # ──────────────────────────────────────────────────────────────
    doc.add_page_break()
    add_heading("4. Dataset Description", 1, NAVY)
    add_body(
        "The project uses 10 curated CSV datasets representing real-world mutual fund "
        "data sourced from AMFI (Association of Mutual Funds in India) and derived from "
        "industry publications."
    )
    doc.add_paragraph()

    tbl = doc.add_table(rows=1, cols=4)
    tbl.style = "Table Grid"
    hdr = tbl.rows[0].cells
    for i, h in enumerate(["File", "Records", "Description", "Key Columns"]):
        hdr[i].text = h
        for run in hdr[i].paragraphs[0].runs:
            run.bold = True
            run.font.size = Pt(9)
    datasets = [
        ("01_fund_master.csv",       "41 rows",    "Fund metadata", "amfi_code, category, benchmark, risk"),
        ("02_nav_history.csv",       "46,000 rows","Daily NAV", "amfi_code, date, nav"),
        ("03_aum_by_fund_house.csv", "90 rows",    "Monthly AUM", "fund_house, date, aum_crore"),
        ("04_monthly_sip_inflows.csv","48 rows",   "Industry SIP", "month, sip_inflow_crore, accounts"),
        ("05_category_inflows.csv",  "144 rows",   "Category flows","category, month, net_inflow"),
        ("06_industry_folio_count.csv","21 rows",  "Folio counts", "equity, debt, hybrid folios"),
        ("07_scheme_performance.csv","41 rows",    "Risk metrics",  "sharpe, alpha, beta, drawdown"),
        ("08_investor_transactions.csv","32,778 rows","Transactions","type, amount, payment_mode"),
        ("09_portfolio_holdings.csv","322 rows",   "Stock holdings","sector, weight_pct"),
        ("10_benchmark_indices.csv", "8,050 rows", "Index prices",  "index_name, date, close"),
    ]
    for row_data in datasets:
        add_table_row(tbl, row_data)

    # ──────────────────────────────────────────────────────────────
    # 5. ETL PIPELINE
    # ──────────────────────────────────────────────────────────────
    doc.add_page_break()
    add_heading("5. ETL Pipeline (D1)", 1, NAVY)
    add_body(
        "The Extract-Transform-Load pipeline standardises all raw data into a validated "
        "processed layer before database insertion. Three scripts handle each stage:"
    )
    doc.add_paragraph()
    etl_steps = [
        ("Extract (data_ingestion.py)",
         "Reads raw CSVs, validates schema completeness, checks amfi_code referential integrity, and fetches live NAV data from mfapi.in API for 5 selected schemes."),
        ("Transform (data_cleaning.py)",
         "Applies pandas forward-fill for missing NAV dates (weekends/holidays), standardises date formats to ISO 8601, normalises portfolio weights to ensure sum <= 100%, and removes invalid transaction records."),
        ("Load (database_loading.py)",
         "Creates 11 SQLite tables per star-schema DDL (schema.sql), validates foreign key references, and bulk-inserts all 10 processed CSVs with zero constraint violations."),
    ]
    for title, desc in etl_steps:
        add_body(f"{title}", bold=True)
        add_body(desc)
        doc.add_paragraph()

    # ──────────────────────────────────────────────────────────────
    # 6. DATABASE DESIGN
    # ──────────────────────────────────────────────────────────────
    doc.add_page_break()
    add_heading("6. Database Design (D2)", 1, NAVY)
    add_body(
        "The SQLite database implements a Star Schema optimised for financial reporting queries. "
        "The schema contains 2 dimension tables and 9 fact/helper tables."
    )
    doc.add_paragraph()
    add_body("Schema Summary:", bold=True)
    schema_rows = [
        ("dim_fund",              "40",     "Fund metadata dimension"),
        ("dim_date",              "1,297",  "Pre-populated date dimension"),
        ("fact_nav",              "46,000", "Daily NAV history"),
        ("fact_performance",      "40",     "Risk-adjusted metrics"),
        ("fact_transactions",     "32,778", "Investor transaction log"),
        ("fact_portfolio",        "322",    "Stock holdings"),
        ("fact_aum",              "90",     "Monthly AUM by fund house"),
        ("fact_benchmark_indices","8,050",  "Index closing prices"),
        ("fact_sip_industry",     "48",     "Monthly SIP inflows"),
        ("fact_category_inflows", "144",    "Net inflows by category"),
        ("fact_industry_folios",  "21",     "Total folio counts"),
    ]
    tbl2 = doc.add_table(rows=1, cols=3)
    tbl2.style = "Table Grid"
    for i, h in enumerate(["Table", "Rows", "Purpose"]):
        tbl2.rows[0].cells[i].text = h
        for run in tbl2.rows[0].cells[i].paragraphs[0].runs:
            run.bold = True
            run.font.size = Pt(9)
    for row in schema_rows:
        add_table_row(tbl2, row)

    # ──────────────────────────────────────────────────────────────
    # 7. EXPLORATORY DATA ANALYSIS
    # ──────────────────────────────────────────────────────────────
    doc.add_page_break()
    add_heading("7. Exploratory Data Analysis (D3)", 1, NAVY)
    eda_insights = [
        ("NAV Trends",
         "Equity large-cap funds show cyclical fluctuations with strong upward growth. "
         "Liquid debt schemes remain highly linear, prioritising wealth preservation.",
         CHART_IMAGES["nav_trends"]),
        ("SIP Inflow Growth",
         "Active SIP accounts grew from ~4.9 Cr to 9.35 Cr; monthly inflows rose from "
         "INR 11,500 Cr to INR 31,002 Cr – demonstrating resilient retail investor confidence.",
         CHART_IMAGES["sip_inflows"]),
        ("AUM Distribution",
         "SBI Mutual Fund leads with INR 12.5 Lakh Crore AUM, followed by ICICI Prudential "
         "(10.74) and HDFC (9.3). Bank-backed fund houses dominate due to branch network leverage.",
         CHART_IMAGES["aum_growth"]),
        ("Risk-Return Scatter",
         "Large-cap and mid-cap equity schemes offer 12–15% CAGR at 13–19% volatility. "
         "Gilt and liquid funds cluster in the low-risk, low-return quadrant.",
         CHART_IMAGES["risk_return"]),
    ]
    for title, desc, img in eda_insights:
        add_heading(title, 2, BLUE)
        add_body(desc)
        add_img(img, width=Inches(5.0), caption=f"Figure: {title}")
        doc.add_paragraph()

    # ──────────────────────────────────────────────────────────────
    # 8. PERFORMANCE ANALYTICS
    # ──────────────────────────────────────────────────────────────
    doc.add_page_break()
    add_heading("8. Performance Analytics (D4)", 1, NAVY)

    formulas = [
        ("CAGR (Compound Annual Growth Rate)",
         "CAGR = (End_NAV / Start_NAV)^(1/years) - 1",
         "Measures annualised NAV growth over a multi-year period, eliminating volatility noise."),
        ("Sharpe Ratio",
         "Sharpe = (R_p - R_f) / sigma_p   [R_f = 5%]",
         "Quantifies excess return per unit of total risk. Funds with Sharpe > 1.0 offer superior risk-adjusted performance."),
        ("Sortino Ratio",
         "Sortino = (R_p - R_f) / downside_sigma",
         "Variant of Sharpe using only downside volatility – more appropriate for asymmetric return distributions."),
        ("Beta",
         "Beta = Cov(R_fund, R_market) / Var(R_market)",
         "Measures market sensitivity. Beta < 1 implies defensive behaviour; Beta > 1 implies amplified market moves."),
        ("Jensen's Alpha",
         "Alpha = R_p - [R_f + Beta * (R_m - R_f)]",
         "Measures excess return above what CAPM predicts. Positive alpha indicates superior fund manager skill."),
        ("Value at Risk (VaR 95%)",
         "Historical: 5th percentile of return distribution | Parametric: mu - 1.645*sigma",
         "Maximum expected 1-day loss at 95% confidence. Equity small-cap funds show VaR ~2% daily; liquid funds near zero."),
        ("Maximum Drawdown",
         "MDD = (Peak_NAV - Trough_NAV) / Peak_NAV",
         "Largest peak-to-trough NAV decline in history. Small-cap funds show MDD of 13–33%; gilt funds < 3%."),
    ]
    for title, formula, interp in formulas:
        add_heading(title, 2, BLUE)
        fb = doc.add_paragraph()
        fr = fb.add_run(f"  {formula}")
        fr.font.name = "Courier New"
        fr.font.size = Pt(10)
        fr.font.color.rgb = RGBColor(0x00, 0x7A, 0xCC)
        add_body(interp)
        doc.add_paragraph()

    add_heading("Top Performing Schemes (Sharpe > 1.0)", 2, BLUE)
    tbl3 = doc.add_table(rows=1, cols=5)
    tbl3.style = "Table Grid"
    for i, h in enumerate(["Scheme", "Category", "3yr CAGR", "Sharpe", "Max DD"]):
        tbl3.rows[0].cells[i].text = h
        for run in tbl3.rows[0].cells[i].paragraphs[0].runs:
            run.bold = True; run.font.size = Pt(9)
    perf_data = [
        ("HDFC Top 100",         "Large Cap", "14.84%", "1.06", "-17.4%"),
        ("Mirae Asset Large Cap","Large Cap", "14.81%", "1.06", "-17.1%"),
        ("ICICI Pru Bluechip D", "Large Cap", "14.41%", "1.03", "-26.6%"),
        ("HDFC Short Term Debt", "Debt",      "7.37%",  "1.84", "-6.0%"),
        ("SBI Magnum Gilt",      "Gilt",      "6.07%",  "1.52", "-2.3%"),
    ]
    for row in perf_data:
        add_table_row(tbl3, row)

    add_img(CHART_IMAGES["var_comparison"], caption="Chart: Historical vs Parametric VaR (95%) by Scheme")
    add_img(CHART_IMAGES["max_drawdown"],   caption="Chart: Average Maximum Drawdown by Category")

    # ──────────────────────────────────────────────────────────────
    # 9. DASHBOARD OVERVIEW
    # ──────────────────────────────────────────────────────────────
    doc.add_page_break()
    add_heading("9. Dashboard Overview (D5)", 1, NAVY)
    add_body(
        "The Streamlit dashboard provides an interactive multi-page interface for "
        "exploring all analytics. Key features include: live NAV filtering by scheme, "
        "Bollinger Band volatility channels, wealth compounding simulators, portfolio "
        "sector overlap analysis, and one-click PDF report export."
    )
    doc.add_paragraph()
    dash_screens = [
        ("Executive Overview & KPIs", CHART_IMAGES["dash_overview"],
         "Landing page showing platform-wide AUM, SIP inflows, active accounts, and scheme count KPI cards."),
        ("Performance Analytics", CHART_IMAGES["dash_perf"],
         "Detailed risk-adjusted metrics table with Sharpe, Sortino, Beta, Alpha, and Drawdown columns."),
        ("Portfolio Allocation", CHART_IMAGES["dash_portfolio"],
         "Sector allocation breakdown across selected funds with overlap detection."),
        ("Risk Analytics", CHART_IMAGES["dash_risk"],
         "Bollinger Band NAV chart, rolling returns, and drawdown timeline for selected scheme."),
    ]
    for title, img, cap in dash_screens:
        add_heading(title, 2, BLUE)
        add_body(cap)
        add_img(img, width=Inches(5.5), caption=f"Dashboard: {title}")
        doc.add_paragraph()

    # ──────────────────────────────────────────────────────────────
    # 10. RECOMMENDATION ENGINE
    # ──────────────────────────────────────────────────────────────
    doc.add_page_break()
    add_heading("10. Recommendation Engine (D6)", 1, NAVY)
    add_body(
        "A rule-based fund recommender (scripts/recommender.py) matches investor profiles "
        "to suitable mutual fund schemes using a weighted composite scoring model."
    )
    doc.add_paragraph()
    add_body("Input Parameters:", bold=True)
    add_bullet("Risk Appetite: Conservative | Moderate | Aggressive")
    add_bullet("Investment Horizon: in years (e.g., 2, 5, 10)")
    add_bullet("Preferred Category: Large Cap, Mid Cap, Small Cap, Gilt, etc.")
    doc.add_paragraph()
    add_body("Scoring Formula:", bold=True)
    fb = doc.add_paragraph()
    fr = fb.add_run("  Score = sum(weight_i * normalised(metric_i))  [scaled 0-100]")
    fr.font.name = "Courier New"; fr.font.size = Pt(10)
    fr.font.color.rgb = RGBColor(0x00, 0x7A, 0xCC)
    doc.add_paragraph()

    tbl4 = doc.add_table(rows=1, cols=4)
    tbl4.style = "Table Grid"
    for i, h in enumerate(["Metric", "Conservative", "Moderate", "Aggressive"]):
        tbl4.rows[0].cells[i].text = h
        for run in tbl4.rows[0].cells[i].paragraphs[0].runs:
            run.bold = True; run.font.size = Pt(9)
    weight_data = [
        ("Sharpe Ratio",          "40%", "30%", "20%"),
        ("3-Year CAGR",           "25%", "35%", "45%"),
        ("Volatility (inverted)", "20%", "15%", "10%"),
        ("Max Drawdown (inv.)",   "15%", "20%", "25%"),
    ]
    for row in weight_data:
        add_table_row(tbl4, row)

    add_img(CHART_IMAGES["rec_scores"], caption="Chart: Recommendation Scores by Risk Profile")

    # ──────────────────────────────────────────────────────────────
    # 11. KEY BUSINESS INSIGHTS
    # ──────────────────────────────────────────────────────────────
    doc.add_page_break()
    add_heading("11. Key Business Insights", 1, NAVY)
    insights = [
        "Small Cap funds deliver highest CAGR (~21-24%) but at ~25% annualised volatility – suitable only for 7+ year horizons.",
        "Direct Plans outperform Regular Plans by 60-120 bps per annum due to lower expense ratios; compounding amplifies this gap significantly over time.",
        "Large Cap and Index funds are highly correlated (rho > 0.9) – holding both provides negligible diversification benefit.",
        "Gilt funds show near-zero or negative correlation with equity during market stress, making them effective portfolio hedges.",
        "Historical VaR exceeds Parametric VaR for equity funds because returns exhibit negative skewness and fat tails – Gaussian assumptions understate tail risk.",
        "SBI, ICICI Prudential, and HDFC Mutual Fund control over 52% of total AUM, reflecting the competitive advantage of bank-distribution networks.",
        "SIP active accounts grew 91% (from 4.9 Cr to 9.35 Cr) in 3 years, signalling a structural shift toward disciplined long-term retail investing.",
        "Flexi-Cap funds offer the best risk-adjusted trade-off (moderate volatility, competitive CAGR, Sharpe > 0.9) for moderate investors.",
        "Financial Services and IT sectors dominate equity fund allocations (30-48% weight), creating hidden concentration risk across seemingly diversified portfolios.",
        "Liquid funds serve as efficient cash-parking vehicles with same-day redemption, but their post-tax returns barely exceed savings account rates.",
    ]
    for i, insight in enumerate(insights, 1):
        add_bullet(f"Insight {i}: {insight}")

    add_img(CHART_IMAGES["risk_return_adv"], caption="Chart: Risk-Return by Fund Category (bubble size = AUM)")
    add_img(CHART_IMAGES["sharpe_violin"],    caption="Chart: Sharpe Ratio Distribution by Equity Category")

    # ──────────────────────────────────────────────────────────────
    # 12. CHALLENGES FACED
    # ──────────────────────────────────────────────────────────────
    doc.add_page_break()
    add_heading("12. Challenges Faced", 1, NAVY)
    challenges = [
        ("Missing NAV Dates",
         "Raw daily NAV datasets contained gaps for weekends and public holidays. "
         "Solution: Forward-fill (ffill) grouped by scheme to carry last known NAV forward accurately."),
        ("SQLite Thread Concurrency",
         "Streamlit's multi-threaded architecture caused 'SQLite objects created in a thread' errors. "
         "Solution: Used @st.cache_resource with check_same_thread=False for safe connection caching."),
        ("Portfolio Weight Anomalies",
         "Certain portfolio holdings had weights summing > 100% or negative values. "
         "Solution: Implemented normalisation and outlier-dropping rules in data_cleaning.py."),
        ("VaR Method Comparison",
         "Parametric VaR requires normal distribution assumption which is violated for equity funds. "
         "Solution: Both Historical and Parametric VaR computed and compared, with explicit explanation of tail-risk underestimation."),
        ("nbformat Cell ID Warnings",
         "Jupyter nbformat 5.1+ requires unique cell IDs. "
         "Solution: Added UUID-based cell IDs to all generated notebook cells."),
        ("Correlation Data Alignment",
         "NAV pivot tables had different date ranges per fund. "
         "Solution: Used dropna() on the returns pivot to ensure only common date ranges are used for correlation."),
    ]
    for title, desc in challenges:
        add_heading(title, 2, BLUE)
        add_body(desc)

    # ──────────────────────────────────────────────────────────────
    # 13. FUTURE IMPROVEMENTS
    # ──────────────────────────────────────────────────────────────
    doc.add_page_break()
    add_heading("13. Future Improvements", 1, NAVY)
    future = [
        "Advanced ML Forecasting: Replace polynomial curve fitting with LSTM or ARIMA models to incorporate macroeconomic indicators (inflation, RBI repo rate) for NAV prediction.",
        "User Authentication & Portfolios: Add login system enabling users to build virtual watchlists, track personalised holdings, and receive SMS/email rebalancing alerts.",
        "Live Data Integration: Connect to NSE/BSE WebSocket feeds for real-time intraday NAV updates and live portfolio re-valuation.",
        "Portfolio Overlap Widget: Integrate a 'Portfolio Overlap Checker' that alerts investors when two selected schemes have > 75% sector overlap.",
        "Mobile-First Dashboard: Develop a React Native or Flutter app for mobile access with push notifications for drawdown alerts.",
        "ESG Scoring: Incorporate Environmental, Social, and Governance (ESG) ratings alongside financial metrics for responsible investing recommendations.",
        "Multi-Asset Expansion: Extend the platform to cover bonds, ETFs, real estate investment trusts (REITs), and international funds.",
        "Automated PDF Scheduling: Implement a cron-based system to auto-generate and email monthly performance PDFs to registered users.",
    ]
    for i, f in enumerate(future, 1):
        add_bullet(f"{i}. {f}")

    # ──────────────────────────────────────────────────────────────
    # 14. CONCLUSION
    # ──────────────────────────────────────────────────────────────
    doc.add_page_break()
    add_heading("14. Conclusion", 1, NAVY)
    add_body(
        "The Bluestock Mutual Fund Analytics Platform successfully demonstrates an "
        "end-to-end data engineering and financial analytics pipeline. Starting from raw "
        "CSV datasets, the project builds a validated relational database, applies "
        "institutional-grade financial metrics, surfaces insights through an interactive "
        "dashboard, and provides personalised investment recommendations through a "
        "rule-based engine."
    )
    doc.add_paragraph()
    add_body(
        "The platform is designed to democratise professional mutual fund analytics for "
        "retail investors, providing the same quality of risk-adjusted performance "
        "evaluation previously available only to institutional players. The modular "
        "architecture ensures easy extension to additional asset classes, real-time data "
        "sources, and machine learning forecasting models."
    )
    doc.add_paragraph()
    add_body(
        "All seven deliverables (D1–D7) have been completed, documented, and validated. "
        "The project represents a production-ready blueprint for a fintech analytics "
        "product serving India's rapidly growing mutual fund ecosystem.",
        bold=False
    )

    # ── Save ──────────────────────────────────────────────────────
    docx_path = REPORTS_DIR / "Final_Report.docx"
    doc.save(str(docx_path))
    print(f"[OK] DOCX saved -> {docx_path}")
    return docx_path


# ─────────────────────────────────────────────────────────────────────────────
# SECTION 2: PDF REPORT (via ReportLab)
# ─────────────────────────────────────────────────────────────────────────────

def build_pdf():
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.lib.units import cm, inch
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, PageBreak,
        Table, TableStyle, Image, HRFlowable, KeepTogether
    )
    from reportlab.platypus.flowables import BalancedColumns
    from reportlab.graphics.shapes import Drawing, Rect
    from reportlab.pdfgen import canvas as cv_mod

    pdf_path = REPORTS_DIR / "Final_Report.pdf"

    # ── Colour definitions ────────────────────────────────────────
    NAVY    = colors.HexColor("#0D1B2A")
    BLUE    = colors.HexColor("#1B6CA8")
    ACCENT  = colors.HexColor("#2ECC71")
    LGREY   = colors.HexColor("#F4F6F9")
    DGREY   = colors.HexColor("#55606B")
    WHITE   = colors.white

    # ── Document ──────────────────────────────────────────────────
    doc = SimpleDocTemplate(
        str(pdf_path),
        pagesize=A4,
        leftMargin=2.5*cm, rightMargin=2.5*cm,
        topMargin=2*cm,    bottomMargin=2*cm,
        title="Bluestock Mutual Fund Analytics – Final Report",
        author="Bluestock Internship Team",
        subject="Capstone Project – Final Documentation",
    )

    styles = getSampleStyleSheet()

    # ── Custom styles ─────────────────────────────────────────────
    def style(name, **kwargs):
        s = ParagraphStyle(name, **kwargs)
        return s

    S_TITLE   = style("MyTitle",   parent=styles["Title"],   fontSize=26, textColor=NAVY,  spaceAfter=6,  alignment=TA_CENTER, fontName="Helvetica-Bold")
    S_SUB     = style("MySub",     parent=styles["Normal"],  fontSize=13, textColor=BLUE,  spaceAfter=4,  alignment=TA_CENTER)
    S_META    = style("MyMeta",    parent=styles["Normal"],  fontSize=10, textColor=DGREY, spaceAfter=2,  alignment=TA_CENTER)
    S_H1      = style("MyH1",      parent=styles["Heading1"],fontSize=16, textColor=NAVY,  spaceBefore=14,spaceAfter=6, fontName="Helvetica-Bold", borderPad=4)
    S_H2      = style("MyH2",      parent=styles["Heading2"],fontSize=12, textColor=BLUE,  spaceBefore=8, spaceAfter=4, fontName="Helvetica-Bold")
    S_BODY    = style("MyBody",    parent=styles["Normal"],  fontSize=10, textColor=colors.black, spaceAfter=6, leading=15, alignment=TA_JUSTIFY)
    S_BULLET  = style("MyBullet",  parent=styles["Normal"],  fontSize=10, leftIndent=16, bulletIndent=4, spaceAfter=3, leading=14)
    S_CAPTION = style("MyCaption", parent=styles["Normal"],  fontSize=8,  textColor=DGREY, alignment=TA_CENTER, spaceAfter=4, fontName="Helvetica-Oblique")
    S_CODE    = style("MyCode",    parent=styles["Code"],    fontSize=9,  textColor=BLUE,  fontName="Courier",  backColor=LGREY, borderPad=6, spaceAfter=6)
    S_TH      = style("MyTH",      parent=styles["Normal"],  fontSize=8,  textColor=WHITE, fontName="Helvetica-Bold")
    S_TD      = style("MyTD",      parent=styles["Normal"],  fontSize=8,  textColor=colors.black)

    def h1(txt):
        return [HRFlowable(width="100%", thickness=2, color=NAVY, spaceAfter=4),
                Paragraph(txt, S_H1)]

    def h2(txt):
        return [Paragraph(txt, S_H2)]

    def body(txt):
        return [Paragraph(txt, S_BODY), Spacer(1, 4)]

    def bullet(txt):
        return [Paragraph(f"• {txt}", S_BULLET)]

    def caption(txt):
        return [Paragraph(txt, S_CAPTION)]

    def img(path, width=14*cm, cap=None):
        p = Path(path)
        items = []
        if p.exists():
            items.append(Image(str(p), width=width, height=width * 0.6))
        else:
            items.append(Paragraph(f"[Chart: {p.name}]", S_CAPTION))
        if cap:
            items += caption(cap)
        items.append(Spacer(1, 6))
        return items

    def make_table(header, rows, col_widths=None):
        data = [[Paragraph(c, S_TH) for c in header]]
        for row in rows:
            data.append([Paragraph(str(c), S_TD) for c in row])
        tbl = Table(data, colWidths=col_widths)
        tbl.setStyle(TableStyle([
            ("BACKGROUND",  (0,0), (-1,0),  NAVY),
            ("ROWBACKGROUNDS", (0,1), (-1,-1), [WHITE, LGREY]),
            ("GRID",        (0,0), (-1,-1), 0.4, colors.HexColor("#CCCCCC")),
            ("VALIGN",      (0,0), (-1,-1), "MIDDLE"),
            ("TOPPADDING",  (0,0), (-1,-1), 4),
            ("BOTTOMPADDING",(0,0),(-1,-1), 4),
            ("LEFTPADDING", (0,0), (-1,-1), 5),
        ]))
        return [tbl, Spacer(1, 8)]

    # ── Build story ───────────────────────────────────────────────
    story = []

    # Title page
    story += [
        Spacer(1, 2*cm),
        Paragraph("Bluestock Mutual Fund Analytics Platform", S_TITLE),
        Spacer(1, 0.3*cm),
        Paragraph("Final Project Report – Capstone Submission", S_SUB),
        Spacer(1, 0.5*cm),
        HRFlowable(width="60%", thickness=2, color=BLUE, hAlign="CENTER"),
        Spacer(1, 0.4*cm),
        Paragraph(f"Date: {TODAY}", S_META),
        Paragraph("Author: Bluestock Internship Team", S_META),
        Paragraph("Institution: Bluestock Fintech Private Limited", S_META),
        Paragraph("Project: Mutual Fund Analytics Capstone (D1 – D7)", S_META),
        PageBreak(),
    ]

    # 1. Executive Summary
    story += h1("1. Executive Summary")
    story += body(
        "The Bluestock Mutual Fund Analytics Platform is an enterprise-grade financial "
        "data engineering and analytics solution. Built as a seven-deliverable capstone "
        "project, the platform ingests 10 raw CSV datasets containing over 46,000 NAV "
        "records, 32,778 investor transactions, and active portfolio holdings. It "
        "orchestrates a full ETL pipeline, populates a relational SQLite star-schema "
        "database, and surfaces advanced financial analytics through an interactive "
        "Streamlit dashboard and a rule-based recommendation engine."
    )
    for m in [
        "Total Platform AUM: INR 62.74 Lakh Crore across 10 fund houses",
        "Monthly Industry SIP Inflows: INR 31,002 Crore (Dec 2025)",
        "Active SIP Accounts: 9.35 Crore | Investor Transactions: 32,778",
        "NAV Records: 46,000+ across 40 schemes (2022–2025)",
        "Database Tables: 11 | Charts Generated: 19+ | Deliverables: D1–D7",
    ]:
        story += bullet(m)
    story.append(Spacer(1, 8))

    # 2. Problem Statement
    story += h1("2. Problem Statement")
    story += body(
        "India's mutual fund industry manages over INR 50 Lakh Crore across thousands "
        "of schemes, yet retail investors lack access to professional-grade risk analytics, "
        "portfolio diversification tools, and personalised investment guidance. Existing "
        "platforms focus on raw returns without communicating risk-adjusted performance, "
        "drawdown exposure, or correlation-driven diversification."
    )
    for p in [
        "Fragmented raw data across multiple CSVs with missing values and inconsistent formats",
        "No centralised relational database for multi-dimensional querying",
        "Retail investors unable to compare funds on Sharpe Ratio, VaR, or drawdown metrics",
        "Absence of a personalised, rule-based fund recommendation engine",
        "No interactive dashboard for real-time exploration of fund performance",
    ]:
        story += bullet(p)

    # 3. Objectives
    story += h1("3. Objectives")
    for d, obj in [
        ("D1 – ETL Pipeline", "Ingest, clean, and standardise 10 raw CSV datasets."),
        ("D2 – SQLite Database", "Design and populate a star-schema relational database."),
        ("D3 – EDA", "Generate 7+ premium static charts revealing key trends."),
        ("D4 – Performance Analytics", "Compute CAGR, Sharpe, Sortino, Beta, Alpha, VaR, Drawdown."),
        ("D5 – Streamlit Dashboard", "Build a multi-page interactive dashboard."),
        ("D6 – Advanced Analytics", "VaR analysis, cohort analysis, correlation, recommendation engine."),
        ("D7 – Documentation", "Professional DOCX + PDF report and PPTX presentation."),
    ]:
        story += bullet(f"<b>{d}:</b> {obj}")

    # 4. Dataset Description
    story += h1("4. Dataset Description")
    story += body("The project uses 10 curated CSV datasets representing real-world mutual fund data from AMFI.")
    story += make_table(
        ["File", "Records", "Description"],
        [
            ("01_fund_master.csv",        "41",     "Fund metadata – category, benchmark, risk"),
            ("02_nav_history.csv",        "46,000", "Daily NAV per scheme (2022–2025)"),
            ("03_aum_by_fund_house.csv",  "90",     "Monthly AUM by fund house"),
            ("04_monthly_sip_inflows.csv","48",     "Industry-wide SIP inflow trends"),
            ("05_category_inflows.csv",   "144",    "Net inflows by category"),
            ("07_scheme_performance.csv", "41",     "Risk metrics: Sharpe, Beta, Alpha, Drawdown"),
            ("08_investor_transactions.csv","32,778","Investor transaction log"),
            ("09_portfolio_holdings.csv", "322",    "Stock holdings and sector weights"),
            ("10_benchmark_indices.csv",  "8,050",  "Market index closing prices"),
        ],
        col_widths=[5.5*cm, 2*cm, 8.5*cm],
    )

    # 5. ETL Pipeline
    story += h1("5. ETL Pipeline (D1)")
    story += body("Three-stage ETL pipeline standardises raw data into a validated processed layer:")
    for stage, desc in [
        ("Extract (data_ingestion.py)",
         "Reads raw CSVs, validates schema completeness, checks referential integrity, and fetches live NAV from mfapi.in API for 5 schemes."),
        ("Transform (data_cleaning.py)",
         "Applies pandas forward-fill for missing NAV dates, standardises date formats to ISO 8601, normalises portfolio weights, and removes invalid records."),
        ("Load (database_loading.py)",
         "Creates 11 SQLite tables per star-schema DDL, validates foreign key references, and bulk-inserts all processed CSVs with zero constraint violations."),
    ]:
        story += h2(stage)
        story += body(desc)

    # 6. Database Design
    story += h1("6. Database Design (D2)")
    story += body("SQLite star-schema with 2 dimension tables and 9 fact/helper tables:")
    story += make_table(
        ["Table", "Rows", "Purpose"],
        [
            ("dim_fund",               "40",     "Fund metadata dimension"),
            ("dim_date",               "1,297",  "Pre-populated date dimension"),
            ("fact_nav",               "46,000", "Daily NAV history"),
            ("fact_performance",       "40",     "Risk-adjusted metrics"),
            ("fact_transactions",      "32,778", "Investor transaction log"),
            ("fact_portfolio",         "322",    "Stock holdings"),
            ("fact_aum",               "90",     "Monthly AUM by fund house"),
            ("fact_benchmark_indices", "8,050",  "Index closing prices"),
        ],
        col_widths=[5.5*cm, 2.5*cm, 8*cm],
    )

    # 7. EDA
    story += h1("7. Exploratory Data Analysis (D3)")
    for title, desc, chart_key in [
        ("NAV Trends",
         "Equity large-cap funds show cyclical growth. Liquid debt schemes remain linear.",
         "nav_trends"),
        ("SIP Inflow Growth",
         "Active SIP accounts grew 91% (4.9 Cr to 9.35 Cr); monthly inflows rose 169% to INR 31,002 Cr.",
         "sip_inflows"),
        ("Risk-Return Scatter",
         "Large/mid-cap funds: 12–15% CAGR at 13–19% volatility. Gilt/liquid cluster in low-risk quadrant.",
         "risk_return"),
    ]:
        story += h2(title)
        story += body(desc)
        story += img(CHART_IMAGES[chart_key], cap=f"Figure: {title}")

    # 8. Performance Analytics
    story += h1("8. Performance Analytics (D4)")
    story += make_table(
        ["Metric", "Formula", "Interpretation"],
        [
            ("CAGR",    "(End/Start)^(1/yr)-1",           "Annualised NAV growth, smoothing volatility"),
            ("Sharpe",  "(Rp-Rf)/sigma  [Rf=5%]",         "Excess return per unit of total risk"),
            ("Sortino", "(Rp-Rf)/downside_sigma",          "Penalises only downside volatility"),
            ("Beta",    "Cov(Rf,Rm)/Var(Rm)",              "Market sensitivity; <1 defensive, >1 amplified"),
            ("Alpha",   "Rp-[Rf+Beta*(Rm-Rf)]",           "Return above CAPM prediction"),
            ("Hist VaR","5th percentile of daily returns", "Max 1-day loss at 95% confidence"),
            ("Max DD",  "(Peak-Trough)/Peak",              "Largest historical peak-to-trough decline"),
        ],
        col_widths=[3*cm, 5*cm, 8*cm],
    )
    story += img(CHART_IMAGES["var_comparison"], cap="Chart: Historical vs Parametric VaR (95%)")
    story += img(CHART_IMAGES["category_cohort"], cap="Chart: Category Performance Cohort Heatmap")

    # 9. Dashboard
    story += h1("9. Dashboard Overview (D5)")
    story += body(
        "Multi-page Streamlit dashboard featuring live NAV filtering, Bollinger Band volatility "
        "channels, wealth compounding simulators, portfolio sector overlap analysis, and one-click PDF export."
    )
    for title, key in [
        ("Executive Overview & KPIs",   "dash_overview"),
        ("Performance Analytics",        "dash_perf"),
        ("Portfolio Allocation",         "dash_portfolio"),
        ("Risk Analytics",               "dash_risk"),
    ]:
        story += h2(title)
        story += img(CHART_IMAGES[key], cap=f"Dashboard: {title}")

    # 10. Recommendation Engine
    story += h1("10. Recommendation Engine (D6)")
    story += body(
        "Rule-based fund recommender matching investor profiles to suitable schemes via "
        "weighted composite scoring. Generates natural-language explanations for each recommendation."
    )
    story += make_table(
        ["Metric", "Conservative", "Moderate", "Aggressive"],
        [
            ("Sharpe Ratio",          "40%", "30%", "20%"),
            ("3-Year CAGR",           "25%", "35%", "45%"),
            ("Volatility (inverted)", "20%", "15%", "10%"),
            ("Max Drawdown (inv.)",   "15%", "20%", "25%"),
        ],
        col_widths=[5*cm, 3*cm, 3*cm, 3*cm],
    )
    story += img(CHART_IMAGES["rec_scores"], cap="Chart: Recommendation Scores by Risk Profile")

    # 11. Business Insights
    story += h1("11. Key Business Insights")
    for i, ins in enumerate([
        "Small Cap funds: highest CAGR (~21-24%) but ~25% volatility – suitable for 7+ year horizons only.",
        "Direct Plans outperform Regular Plans by 60-120 bps p.a. – compounding amplifies this gap significantly.",
        "Large Cap and Index funds are redundant in a portfolio (rho > 0.9 correlation).",
        "Gilt funds provide genuine portfolio hedging (near-zero to negative correlation with equity).",
        "Historical VaR > Parametric VaR for equity due to fat tails – Gaussian models understate risk.",
        "SBI, ICICI, HDFC control 52%+ of AUM via bank-distribution network advantages.",
        "SIP accounts grew 91% in 3 years, signalling structural shift to disciplined long-term investing.",
        "Flexi-Cap is the optimal anchor holding for moderate investors (Sharpe > 0.9, moderate volatility).",
        "Banking/IT sector concentration (30-48%) creates hidden correlated risk across 'diversified' portfolios.",
        "Liquid funds are cash-parking tools, not investments – post-tax returns barely exceed savings accounts.",
    ], 1):
        story += bullet(f"<b>#{i}:</b> {ins}")

    # 12. Challenges
    story += h1("12. Challenges Faced")
    for ch, sol in [
        ("Missing NAV dates (weekends/holidays)", "Pandas forward-fill (ffill) grouped by scheme"),
        ("SQLite thread concurrency in Streamlit", "@st.cache_resource + check_same_thread=False"),
        ("Portfolio weight anomalies (sum > 100%)", "Normalisation and outlier-dropping in data_cleaning.py"),
        ("Parametric VaR vs fat-tailed distributions", "Both methods computed; fat-tail explanation added"),
        ("nbformat cell ID warnings",               "UUID-based cell IDs added to all generated cells"),
        ("NAV date-range misalignment for correlations", "dropna() on returns pivot for common date ranges"),
    ]:
        story += h2(ch)
        story += body(f"Solution: {sol}")

    # 13. Future Improvements
    story += h1("13. Future Improvements")
    for f in [
        "LSTM/ARIMA models for NAV forecasting with macro-economic indicators",
        "User authentication with personalised virtual portfolios and rebalancing alerts",
        "Live WebSocket integration for real-time market data feeds",
        "Portfolio Overlap Widget alerting when two funds share > 75% sector allocation",
        "Mobile-first React Native/Flutter app with push notifications",
        "ESG scoring integration for responsible investing recommendations",
        "Multi-asset expansion: bonds, ETFs, REITs, and international funds",
        "Automated monthly PDF reporting via cron-scheduled email delivery",
    ]:
        story += bullet(f)

    # 14. Conclusion
    story += h1("14. Conclusion")
    story += body(
        "The Bluestock Mutual Fund Analytics Platform successfully demonstrates a complete "
        "end-to-end data engineering and financial analytics pipeline – from raw CSV ingestion "
        "to institutional-grade risk metrics, interactive visualisation, and personalised "
        "investment recommendations."
    )
    story += body(
        "All seven deliverables (D1–D7) have been completed, documented, and validated. "
        "The platform represents a production-ready blueprint for a fintech analytics product "
        "serving India's rapidly growing mutual fund ecosystem."
    )

    # ── Build ─────────────────────────────────────────────────────
    doc.build(story)
    print(f"[OK] PDF saved -> {pdf_path}")
    return pdf_path


# ─────────────────────────────────────────────────────────────────────────────
# SECTION 3: POWERPOINT PRESENTATION
# ─────────────────────────────────────────────────────────────────────────────

def build_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt, Cm
    import pptx.oxml.ns as ns
    from lxml import etree

    pptx_path = REPORTS_DIR / "Presentation.pptx"

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ── Colour palette ─────────────────────────────────────────────
    C_NAVY   = RGBColor(0x0D, 0x1B, 0x2A)
    C_BLUE   = RGBColor(0x1B, 0x6C, 0xA8)
    C_ACCENT = RGBColor(0x2E, 0xCC, 0x71)
    C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    C_LGREY  = RGBColor(0xF4, 0xF6, 0xF9)
    C_DGREY  = RGBColor(0x55, 0x60, 0x6B)
    C_GOLD   = RGBColor(0xF3, 0x9C, 0x12)

    SL = prs.slide_layouts

    # ── Helper functions ──────────────────────────────────────────
    def add_rect(slide, left, top, width, height, fill_color, opacity=1.0):
        shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        return shape

    def add_textbox(slide, text, left, top, width, height,
                    font_size=18, bold=False, color=C_WHITE,
                    align=PP_ALIGN.LEFT, italic=False, word_wrap=True):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = word_wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        return txBox

    def add_img_safe(slide, path, left, top, width, height=None):
        p = Path(path)
        if not p.exists():
            return
        if height:
            slide.shapes.add_picture(str(p), Inches(left), Inches(top), Inches(width), Inches(height))
        else:
            slide.shapes.add_picture(str(p), Inches(left), Inches(top), Inches(width))

    def new_slide(layout_idx=6):
        layout = prs.slide_layouts[layout_idx]
        return prs.slides.add_slide(layout)

    def draw_header_bar(slide, title, subtitle=None):
        add_rect(slide, 0, 0, 13.33, 1.3, C_NAVY)
        add_textbox(slide, title, 0.3, 0.1, 10, 0.7, font_size=28, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
        if subtitle:
            add_textbox(slide, subtitle, 0.3, 0.75, 10, 0.45, font_size=13, color=C_LGREY, align=PP_ALIGN.LEFT)
        add_rect(slide, 0, 1.3, 13.33, 0.06, C_BLUE)

    def add_bullet_box(slide, bullets, left, top, width, height,
                       font_size=13, color=C_NAVY, bg=None):
        if bg:
            add_rect(slide, left, top, width, height, bg)
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(3)
            run = p.add_run()
            run.text = f"• {b}"
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
        return txBox

    def draw_footer(slide, slide_num, total=13):
        add_rect(slide, 0, 7.2, 13.33, 0.3, C_NAVY)
        add_textbox(slide, "Bluestock Mutual Fund Analytics Platform  |  Capstone Project",
                    0.2, 7.22, 9, 0.25, font_size=8, color=C_LGREY)
        add_textbox(slide, f"Slide {slide_num} / {total}", 11.5, 7.22, 1.5, 0.25,
                    font_size=8, color=C_LGREY, align=PP_ALIGN.RIGHT)

    def kpi_box(slide, value, label, left, top, width=2.6, height=1.1,
                value_color=C_ACCENT, bg=C_NAVY):
        add_rect(slide, left, top, width, height, bg)
        add_textbox(slide, value, left+0.1, top+0.05, width-0.2, 0.55,
                    font_size=22, bold=True, color=value_color, align=PP_ALIGN.CENTER)
        add_textbox(slide, label, left+0.1, top+0.57, width-0.2, 0.45,
                    font_size=10, color=C_LGREY, align=PP_ALIGN.CENTER)

    TOTAL_SLIDES = 13

    # ══════════════════════════════════════════════════════════════
    # SLIDE 1 – TITLE
    # ══════════════════════════════════════════════════════════════
    sl = new_slide(6)
    add_rect(sl, 0, 0, 13.33, 7.5, C_NAVY)
    add_rect(sl, 0, 5.8, 13.33, 1.7, C_BLUE)

    add_textbox(sl, "Bluestock", 0.8, 0.8, 11.5, 1.2,
                font_size=52, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, "Mutual Fund Analytics Platform", 0.8, 1.9, 11.5, 0.8,
                font_size=28, bold=False, color=C_LGREY, align=PP_ALIGN.CENTER)
    add_rect(sl, 3, 2.85, 7.33, 0.06, C_ACCENT)
    add_textbox(sl, "Final Capstone Project  |  Deliverables D1 – D7", 0.8, 3.0, 11.5, 0.6,
                font_size=16, color=C_LGREY, align=PP_ALIGN.CENTER)
    add_textbox(sl, f"{TODAY}  |  Bluestock Fintech Private Limited", 0.8, 3.65, 11.5, 0.5,
                font_size=13, color=C_DGREY, align=PP_ALIGN.CENTER, italic=True)

    kpi_box(sl, "INR 62.74L Cr", "Platform AUM",      0.5,  5.95, value_color=C_ACCENT, bg=C_BLUE)
    kpi_box(sl, "46,000+",        "NAV Records",        3.37, 5.95, value_color=C_GOLD,   bg=C_BLUE)
    kpi_box(sl, "32,778",         "Transactions",       6.24, 5.95, value_color=C_WHITE,  bg=C_BLUE)
    kpi_box(sl, "D1 – D7",        "Deliverables",       9.11, 5.95, value_color=C_ACCENT, bg=C_BLUE)
    draw_footer(sl, 1, TOTAL_SLIDES)

    # ══════════════════════════════════════════════════════════════
    # SLIDE 2 – PROBLEM STATEMENT
    # ══════════════════════════════════════════════════════════════
    sl = new_slide(6)
    add_rect(sl, 0, 0, 13.33, 7.5, C_LGREY)
    draw_header_bar(sl, "Problem Statement", "Why does this platform matter?")

    add_textbox(sl, "India's mutual fund industry manages INR 50+ Lakh Crore, yet retail investors "
                    "lack professional-grade tools for risk analysis, portfolio diversification, "
                    "and personalised guidance.", 0.4, 1.45, 12.5, 1.0,
                font_size=13, color=C_NAVY)

    problems = [
        "Fragmented data across CSVs — no single source of truth",
        "No risk-adjusted metrics (Sharpe, VaR, Drawdown) for retail investors",
        "Absence of personalised fund recommendation tools",
        "No interactive dashboard for exploration and comparison",
        "Missing portfolio correlation and diversification analytics",
    ]
    add_bullet_box(sl, problems, 0.4, 2.5, 7.5, 4.2, font_size=14, color=C_NAVY)

    add_rect(sl, 8.2, 2.0, 4.7, 4.9, C_NAVY)
    add_textbox(sl, "Our Solution", 8.4, 2.1, 4.3, 0.5, font_size=14, bold=True, color=C_ACCENT)
    solution_pts = ["End-to-end ETL pipeline", "Relational SQLite DB", "7 EDA charts",
                    "Risk analytics engine", "Interactive dashboard", "Recommender system"]
    for i, s in enumerate(solution_pts):
        add_textbox(sl, f"+ {s}", 8.4, 2.65 + i*0.6, 4.2, 0.55, font_size=12, color=C_WHITE)

    draw_footer(sl, 2, TOTAL_SLIDES)

    # ══════════════════════════════════════════════════════════════
    # SLIDE 3 – DATASET
    # ══════════════════════════════════════════════════════════════
    sl = new_slide(6)
    add_rect(sl, 0, 0, 13.33, 7.5, C_LGREY)
    draw_header_bar(sl, "Dataset Description", "10 curated AMFI-sourced CSV datasets")

    datasets_info = [
        ("01 Fund Master",     "41 schemes, categories, risk grades, fund managers"),
        ("02 NAV History",     "46,000+ daily NAV records, 2022–2025"),
        ("03 AUM by AMC",      "90 monthly AUM records by fund house"),
        ("04 SIP Inflows",     "48 months of industry-wide SIP trends"),
        ("05 Category Flows",  "144 net inflow records by category"),
        ("07 Performance",     "41 schemes with Sharpe, Beta, Alpha, Drawdown"),
        ("08 Transactions",    "32,778 investor transaction records"),
        ("09 Portfolio",       "322 stock holdings with sector weights"),
        ("10 Benchmarks",      "8,050 index closing prices (Nifty 50, 100, etc.)"),
    ]
    for i, (name, desc) in enumerate(datasets_info):
        row, col = divmod(i, 3)
        bx_left = 0.3 + col * 4.35
        bx_top  = 1.6 + row * 1.75
        add_rect(sl, bx_left, bx_top, 4.1, 1.55, C_NAVY)
        add_textbox(sl, name, bx_left+0.1, bx_top+0.05, 3.9, 0.45,
                    font_size=12, bold=True, color=C_ACCENT)
        add_textbox(sl, desc, bx_left+0.1, bx_top+0.5, 3.9, 0.95,
                    font_size=10, color=C_LGREY, word_wrap=True)
    draw_footer(sl, 3, TOTAL_SLIDES)

    # ══════════════════════════════════════════════════════════════
    # SLIDE 4 – ETL PIPELINE
    # ══════════════════════════════════════════════════════════════
    sl = new_slide(6)
    add_rect(sl, 0, 0, 13.33, 7.5, C_LGREY)
    draw_header_bar(sl, "ETL Pipeline (D1)", "Extract → Transform → Load")

    stages = [
        ("EXTRACT", "data_ingestion.py",
         ["Reads 10 raw CSVs", "Validates AMFI code integrity", "Fetches live NAV from mfapi.in API"]),
        ("TRANSFORM", "data_cleaning.py",
         ["Forward-fills missing NAV dates", "Standardises date formats to ISO 8601", "Normalises portfolio weights"]),
        ("LOAD", "database_loading.py",
         ["Creates 11 SQLite tables", "Validates foreign key constraints", "Zero violations confirmed"]),
    ]
    arrow_colors = [C_ACCENT, C_GOLD, C_BLUE]
    for i, (stage, script, pts) in enumerate(stages):
        bx = 0.3 + i * 4.3
        add_rect(sl, bx, 1.6, 4.05, 5.2, arrow_colors[i])
        add_textbox(sl, stage,  bx+0.1, 1.65, 3.85, 0.55, font_size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(sl, script, bx+0.1, 2.2,  3.85, 0.45, font_size=10, color=C_NAVY,  align=PP_ALIGN.CENTER, italic=True)
        add_rect(sl, bx+0.1, 2.7, 3.85, 0.05, C_WHITE)
        for j, pt in enumerate(pts):
            add_textbox(sl, f"• {pt}", bx+0.15, 2.8 + j*0.9, 3.75, 0.85, font_size=12, color=C_WHITE)
        if i < 2:
            add_textbox(sl, "→", bx+4.1, 3.5, 0.3, 0.5, font_size=28, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

    kpi_box(sl, "46,002", "Records Loaded",  0.5,  6.85, value_color=C_ACCENT, bg=C_NAVY)
    kpi_box(sl, "11",     "DB Tables",        3.37, 6.85, value_color=C_GOLD,   bg=C_NAVY)
    kpi_box(sl, "0",      "FK Violations",    6.24, 6.85, value_color=C_ACCENT, bg=C_NAVY)
    kpi_box(sl, "100%",   "Completeness",     9.11, 6.85, value_color=C_WHITE,  bg=C_NAVY)
    draw_footer(sl, 4, TOTAL_SLIDES)

    # ══════════════════════════════════════════════════════════════
    # SLIDE 5 – DATABASE
    # ══════════════════════════════════════════════════════════════
    sl = new_slide(6)
    add_rect(sl, 0, 0, 13.33, 7.5, C_LGREY)
    draw_header_bar(sl, "Database Design (D2)", "SQLite Star Schema – 11 Tables")

    add_textbox(sl, "DIMENSION TABLES", 0.3, 1.45, 6.1, 0.4, font_size=11, bold=True, color=C_BLUE)
    for name, rows in [("dim_fund", "40 rows"), ("dim_date", "1,297 rows")]:
        r = (0.3 if name == "dim_fund" else 3.5)
        add_rect(sl, r, 1.9, 2.9, 0.8, C_BLUE)
        add_textbox(sl, f"{name}\n{rows}", r+0.1, 1.95, 2.7, 0.7, font_size=12, color=C_WHITE)

    add_textbox(sl, "FACT TABLES", 0.3, 2.85, 12.5, 0.4, font_size=11, bold=True, color=C_NAVY)
    facts = [
        ("fact_nav",          "46,000"),
        ("fact_performance",  "40"),
        ("fact_transactions", "32,778"),
        ("fact_portfolio",    "322"),
        ("fact_aum",          "90"),
        ("fact_benchmarks",   "8,050"),
        ("fact_sip_industry", "48"),
        ("fact_cat_inflows",  "144"),
        ("fact_folios",       "21"),
    ]
    for i, (name, rows) in enumerate(facts):
        col, row_ = divmod(i, 3)
        bx = 0.3 + col * 4.35
        by = 3.3 + row_ * 1.2
        add_rect(sl, bx, by, 4.1, 1.1, C_NAVY)
        add_textbox(sl, name, bx+0.1, by+0.05, 3.9, 0.45, font_size=11, bold=True, color=C_ACCENT)
        add_textbox(sl, rows + " rows", bx+0.1, by+0.52, 3.9, 0.45, font_size=10, color=C_LGREY)
    draw_footer(sl, 5, TOTAL_SLIDES)

    # ══════════════════════════════════════════════════════════════
    # SLIDE 6 – EDA
    # ══════════════════════════════════════════════════════════════
    sl = new_slide(6)
    add_rect(sl, 0, 0, 13.33, 7.5, C_LGREY)
    draw_header_bar(sl, "Exploratory Data Analysis (D3)", "7 premium charts revealing key market trends")

    add_img_safe(sl, CHART_IMAGES["nav_trends"],  0.2, 1.5, 6.5, 2.7)
    add_img_safe(sl, CHART_IMAGES["risk_return"], 6.8, 1.5, 6.3, 2.7)
    add_textbox(sl, "NAV Trends: Equity vs Debt",
                0.2, 4.3, 6.4, 0.5, font_size=10, color=C_DGREY, italic=True)
    add_textbox(sl, "Risk vs Return Scatter (3-Year)",
                6.8, 4.3, 6.3, 0.5, font_size=10, color=C_DGREY, italic=True)

    obs = [
        "Equity Large-Cap: 12–15% CAGR | Volatility 13–19%",
        "SIP accounts grew 91% in 3 years (4.9 Cr → 9.35 Cr)",
        "Financial Services & IT dominate sector allocation (30-48%)",
        "Liquid/Gilt funds: < 5% volatility, 5–8% returns",
    ]
    add_bullet_box(sl, obs, 0.3, 4.85, 12.5, 2.3, font_size=12, color=C_NAVY)
    draw_footer(sl, 6, TOTAL_SLIDES)

    # ══════════════════════════════════════════════════════════════
    # SLIDE 7 – PERFORMANCE METRICS
    # ══════════════════════════════════════════════════════════════
    sl = new_slide(6)
    add_rect(sl, 0, 0, 13.33, 7.5, C_LGREY)
    draw_header_bar(sl, "Performance Analytics (D4)", "Institutional-grade risk-adjusted metrics for 40 schemes")

    metrics_info = [
        ("CAGR",        "Annualised NAV growth",         "Small Cap: 21-24% | Large Cap: 12-15%"),
        ("Sharpe",      "Excess return / total risk",    "Best: HDFC Top 100 & Mirae (1.06)"),
        ("Sortino",     "Excess return / downside risk", "Gilt funds excel (2.1-2.8)"),
        ("Beta",        "Market sensitivity",            "Liquid funds: 0.1-0.5 | Equity: 0.8-1.1"),
        ("Alpha",       "Skill above CAPM expectation",  "Best: HDFC Short Term (1.98%)"),
        ("VaR (95%)",   "Max 1-day loss at 95% conf.",   "Small Cap: ~2% | Liquid: < 0.05%"),
        ("Max Drawdown","Largest historical decline",    "Small Cap: -13 to -33% | Gilt: < -3%"),
    ]
    for i, (name, def_, val) in enumerate(metrics_info):
        col, row_ = divmod(i, 4)
        bx = 0.3 + col * 6.5
        by = 1.6 + row_ * 1.35
        add_rect(sl, bx, by, 6.2, 1.25, C_NAVY if col == 0 else C_BLUE)
        add_textbox(sl, name, bx+0.1, by+0.05, 5.9, 0.4, font_size=14, bold=True, color=C_GOLD)
        add_textbox(sl, def_, bx+0.1, by+0.43, 5.9, 0.35, font_size=10, color=C_LGREY)
        add_textbox(sl, val,  bx+0.1, by+0.77, 5.9, 0.4,  font_size=10, color=C_ACCENT, bold=True)
    draw_footer(sl, 7, TOTAL_SLIDES)

    # ══════════════════════════════════════════════════════════════
    # SLIDE 8 – DASHBOARD
    # ══════════════════════════════════════════════════════════════
    sl = new_slide(6)
    add_rect(sl, 0, 0, 13.33, 7.5, C_LGREY)
    draw_header_bar(sl, "Interactive Dashboard (D5)", "Multi-page Streamlit app with live filtering and analytics")

    # 2x2 screenshot grid
    screens = [
        (CHART_IMAGES["dash_overview"],  "Executive Overview & KPIs",     0.2,  1.5),
        (CHART_IMAGES["dash_perf"],      "Performance Analytics",          6.9,  1.5),
        (CHART_IMAGES["dash_portfolio"], "Portfolio Allocation",           0.2,  4.5),
        (CHART_IMAGES["dash_risk"],      "Risk Analytics (Bollinger)",     6.9,  4.5),
    ]
    for path, label, left, top in screens:
        add_img_safe(sl, path, left, top, 6.3, 2.7)
        add_textbox(sl, label, left, top + 2.75, 6.3, 0.4,
                    font_size=9, color=C_DGREY, italic=True, align=PP_ALIGN.CENTER)
    draw_footer(sl, 8, TOTAL_SLIDES)

    # ══════════════════════════════════════════════════════════════
    # SLIDE 9 – RECOMMENDATION ENGINE
    # ══════════════════════════════════════════════════════════════
    sl = new_slide(6)
    add_rect(sl, 0, 0, 13.33, 7.5, C_LGREY)
    draw_header_bar(sl, "Recommendation Engine (D6)", "Rule-based fund matching by risk profile")

    # Pipeline flow
    flow = ["INPUT\nRisk + Horizon + Category", "FILTER\nRisk grade & horizon checks", "SCORE\nWeighted composite 0–100", "RANK & EXPLAIN\nTop-N with NL explanation"]
    flow_colors = [C_BLUE, C_GOLD, C_ACCENT, C_NAVY]
    for i, (txt, clr) in enumerate(zip(flow, flow_colors)):
        add_rect(sl, 0.3 + i * 3.1, 1.6, 2.85, 1.4, clr)
        add_textbox(sl, txt, 0.35 + i * 3.1, 1.65, 2.75, 1.3,
                    font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        if i < 3:
            add_textbox(sl, "→", 3.15 + i * 3.1, 2.1, 0.3, 0.5,
                        font_size=22, bold=True, color=C_NAVY)

    # Weight table
    add_textbox(sl, "Scoring Weights by Risk Profile", 0.3, 3.2, 6.5, 0.45,
                font_size=13, bold=True, color=C_NAVY)
    weight_tbl = [
        ["Metric",            "Conservative", "Moderate", "Aggressive"],
        ["Sharpe Ratio",      "40%",          "30%",      "20%"],
        ["3-Year CAGR",       "25%",          "35%",      "45%"],
        ["Volatility (inv.)", "20%",          "15%",      "10%"],
        ["Max Drawdown (inv.)","15%",          "20%",      "25%"],
    ]
    row_colors = [C_NAVY, C_BLUE, C_BLUE, C_BLUE, C_BLUE]
    for ri, row in enumerate(weight_tbl):
        for ci, cell in enumerate(row):
            cx = 0.3 + ci * 1.55
            cy = 3.7 + ri * 0.65
            add_rect(sl, cx, cy, 1.5, 0.62, row_colors[ri] if ri == 0 else (C_LGREY if ri % 2 else C_WHITE))
            fc = C_WHITE if ri == 0 else C_NAVY
            fb = ri == 0
            add_textbox(sl, cell, cx+0.05, cy+0.08, 1.4, 0.48, font_size=10, bold=fb, color=fc, align=PP_ALIGN.CENTER)

    add_img_safe(sl, CHART_IMAGES["rec_scores"], 6.8, 3.0, 6.2, 4.0)
    add_textbox(sl, "27 recommendations generated across 5 investor profiles",
                6.8, 7.05, 6.2, 0.35, font_size=9, color=C_DGREY, italic=True)
    draw_footer(sl, 9, TOTAL_SLIDES)

    # ══════════════════════════════════════════════════════════════
    # SLIDE 10 – KEY INSIGHTS
    # ══════════════════════════════════════════════════════════════
    sl = new_slide(6)
    add_rect(sl, 0, 0, 13.33, 7.5, C_LGREY)
    draw_header_bar(sl, "Key Business Insights", "10 data-driven observations from the analytics pipeline")

    insights = [
        ("Small Cap Risk-Return", "Highest CAGR (21-24%) but ~25% volatility — 7+ year horizon required"),
        ("Direct vs Regular Plans","Direct outperforms by 60-120 bps p.a. — compounding amplifies gap"),
        ("Large Cap / Index Overlap","rho > 0.9 — holding both provides minimal diversification"),
        ("Gilt as Hedge",          "Near-zero to negative correlation with equity during stress periods"),
        ("Fat Tails in Equity",    "Hist. VaR > Parametric VaR — Gaussian models understate risk"),
        ("AMC Concentration",      "SBI + ICICI + HDFC control 52%+ of platform AUM"),
        ("SIP Momentum",           "Active SIP accounts grew 91% in 3 years (4.9 Cr → 9.35 Cr)"),
        ("Flexi-Cap Optimal",      "Best risk-adjusted anchor holding for moderate investors (Sharpe >0.9)"),
        ("Sector Concentration",   "Banking/IT: 30-48% weight creates hidden correlated risk"),
        ("Liquid Fund Reality",    "Post-tax returns barely beat savings accounts — cash-parking only"),
    ]
    for i, (title, body_) in enumerate(insights):
        col, row_ = divmod(i, 5)
        bx = 0.2 + col * 6.6
        by = 1.55 + row_ * 1.12
        bg = C_NAVY if col == 0 else C_BLUE
        add_rect(sl, bx, by, 6.35, 1.05, bg)
        add_textbox(sl, title, bx+0.1, by+0.04, 6.15, 0.38, font_size=11, bold=True, color=C_GOLD)
        add_textbox(sl, body_, bx+0.1, by+0.44, 6.15, 0.55, font_size=9,  color=C_WHITE, word_wrap=True)
    draw_footer(sl, 10, TOTAL_SLIDES)

    # ══════════════════════════════════════════════════════════════
    # SLIDE 11 – CHALLENGES
    # ══════════════════════════════════════════════════════════════
    sl = new_slide(6)
    add_rect(sl, 0, 0, 13.33, 7.5, C_LGREY)
    draw_header_bar(sl, "Challenges & Solutions", "Engineering challenges encountered and resolved")

    challenges = [
        ("Missing NAV Dates",        "Weekends/holidays created gaps in daily NAV",        "pandas ffill() grouped by scheme"),
        ("SQLite Thread Concurrency","Streamlit multi-threading caused DB errors",          "@st.cache_resource + check_same_thread=False"),
        ("Weight Anomalies",         "Portfolio weights summed > 100% or were negative",   "Normalisation & outlier-dropping rules"),
        ("VaR Assumption Mismatch",  "Equity returns violate Gaussian distribution",       "Compute both methods; document fat-tail risk"),
        ("nbformat Cell IDs",        "Jupyter 5.1+ requires unique IDs per cell",          "UUID-based IDs added to all notebook cells"),
        ("Correlation Alignment",    "NAV pivots had different date ranges per fund",      "dropna() on returns pivot for common ranges"),
    ]
    for i, (prob, context, sol) in enumerate(challenges):
        row_, col = divmod(i, 2)
        bx = 0.2 + col * 6.55
        by = 1.55 + row_ * 1.85
        add_rect(sl, bx,      by,      6.3,  1.75, C_NAVY)
        add_rect(sl, bx+0.08, by+0.08, 5.6,  0.45, C_GOLD)
        add_textbox(sl, prob,    bx+0.15, by+0.1,  5.5, 0.4,  font_size=12, bold=True, color=C_NAVY)
        add_textbox(sl, context, bx+0.12, by+0.6,  6.0, 0.5,  font_size=10, color=C_LGREY)
        add_textbox(sl, f"Solution: {sol}", bx+0.12, by+1.12, 6.0, 0.55, font_size=10, color=C_ACCENT, bold=True)
    draw_footer(sl, 11, TOTAL_SLIDES)

    # ══════════════════════════════════════════════════════════════
    # SLIDE 12 – FUTURE SCOPE
    # ══════════════════════════════════════════════════════════════
    sl = new_slide(6)
    add_rect(sl, 0, 0, 13.33, 7.5, C_LGREY)
    draw_header_bar(sl, "Future Scope", "Enhancements planned for production release")

    future_items = [
        ("AI/ML Forecasting",      "LSTM/ARIMA models with macroeconomic indicators (inflation, RBI repo rate)"),
        ("User Portfolios",         "Authentication + virtual watchlists + rebalancing alerts via SMS/email"),
        ("Live Market Data",        "NSE/BSE WebSocket feeds for real-time NAV and portfolio re-valuation"),
        ("Portfolio Overlap Tool",  "Alert when two selected funds share > 75% sector allocation"),
        ("Mobile App",              "React Native/Flutter mobile app with push notifications"),
        ("ESG Integration",         "Environmental, Social, Governance ratings alongside financial metrics"),
        ("Multi-Asset Platform",    "Bonds, ETFs, REITs, and international funds coverage"),
        ("Automated Reporting",     "Cron-scheduled monthly PDF reports emailed to registered users"),
    ]
    for i, (title, desc) in enumerate(future_items):
        col, row_ = divmod(i, 4)
        bx = 0.2 + col * 6.55
        by = 1.55 + row_ * 1.35
        add_rect(sl, bx, by, 6.3, 1.25, C_BLUE if col == 0 else C_NAVY)
        add_textbox(sl, f"+ {title}", bx+0.1, by+0.05, 6.1, 0.45, font_size=12, bold=True, color=C_GOLD)
        add_textbox(sl, desc, bx+0.1, by+0.52, 6.1, 0.68, font_size=10, color=C_WHITE, word_wrap=True)
    draw_footer(sl, 12, TOTAL_SLIDES)

    # ══════════════════════════════════════════════════════════════
    # SLIDE 13 – CONCLUSION
    # ══════════════════════════════════════════════════════════════
    sl = new_slide(6)
    add_rect(sl, 0, 0, 13.33, 7.5, C_NAVY)
    add_rect(sl, 0, 6.8, 13.33, 0.7, C_BLUE)

    add_textbox(sl, "Conclusion", 0.5, 0.4, 12, 0.8, font_size=36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, 3, 1.3, 7.33, 0.06, C_ACCENT)

    add_textbox(sl,
        "The Bluestock Mutual Fund Analytics Platform delivers a complete, "
        "production-ready data engineering and financial analytics solution — "
        "from raw CSV ingestion to institutional-grade risk metrics, "
        "interactive visualisation, and personalised investment recommendations.",
        0.5, 1.5, 12.3, 1.2, font_size=14, color=C_LGREY, align=PP_ALIGN.CENTER)

    conclusion_pts = [
        ("D1–D2", "ETL + Star-Schema Database"),
        ("D3–D4", "EDA + Performance Analytics"),
        ("D5",    "Interactive Streamlit Dashboard"),
        ("D6",    "VaR + Cohort + Recommender"),
        ("D7",    "Final Report & Presentation"),
    ]
    box_colors = [C_BLUE, C_BLUE, C_ACCENT, C_GOLD, C_BLUE]
    for i, ((d, desc), clr) in enumerate(zip(conclusion_pts, box_colors)):
        bx = 0.4 + i * 2.5
        add_rect(sl, bx, 2.9, 2.35, 1.5, clr)
        add_textbox(sl, d,    bx+0.1, 2.95, 2.15, 0.55, font_size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(sl, desc, bx+0.1, 3.5,  2.15, 0.85, font_size=10, color=C_WHITE, align=PP_ALIGN.CENTER, word_wrap=True)

    add_textbox(sl, "All 7 deliverables completed, validated, and documented.",
                0.5, 4.65, 12.3, 0.6, font_size=15, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    add_textbox(sl, "Thank You", 0, 5.3, 13.33, 1.0,
                font_size=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, "Bluestock Fintech Private Limited  |  github.com/eashita2409/bluestock-internship",
                0, 6.82, 13.33, 0.45, font_size=10, color=C_LGREY, align=PP_ALIGN.CENTER)

    prs.save(str(pptx_path))
    print(f"[OK] PPTX saved -> {pptx_path}")
    return pptx_path


# ─────────────────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────────────────

def main():
    print("Generating D7 Final Documentation...")
    print(f"  Output directory: {REPORTS_DIR}")
    print()

    docx_path = build_docx()
    print()
    pdf_path  = build_pdf()
    print()
    pptx_path = build_pptx()
    print()

    print("=" * 60)
    print("D7 Final Documentation Complete")
    print("=" * 60)
    for f in [docx_path, pdf_path, pptx_path]:
        size_kb = Path(f).stat().st_size // 1024
        print(f"  {Path(f).name:<25}  {size_kb:>5} KB")


if __name__ == "__main__":
    main()
